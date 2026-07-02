"""
Asset export helpers.

Two responsibilities:
1. Persist AI-generated images/videos to local disk (static/generated/) so
   they have a stable, downloadable URL instead of just hotlinking a
   third-party endpoint that may expire or rate-limit.
2. Build *real* downloadable files from the structured slide/poster data the
   agent emits: an actual .pptx (via python-pptx) and an actual poster .pdf
   (via reportlab if available, else a pure-Python fallback).

Both are used by new routes in main.py: /export/pptx/{thread_id} and
/export/poster-pdf/{thread_id}, plus /download/asset for saving generated
media.
"""
import os
import re
import uuid
import requests

from config import settings
from logger import get_logger

logger = get_logger("asset_export")

os.makedirs(settings.ASSET_STORAGE_DIR, exist_ok=True)
os.makedirs(settings.EXPORT_STORAGE_DIR, exist_ok=True)


# --------------------------------------------------------------------------- #
# Persist a remote generated asset (image or video) to local disk
# --------------------------------------------------------------------------- #

def save_remote_asset(remote_url: str, kind: str = "image", headers: dict = None) -> dict:
    """Download a remote generated asset and save it locally so it has a
    stable, downloadable URL. Returns {"url": ..., "filename": ...} or raises."""
    ext = ".mp4" if kind == "video" else ".png"
    if "." in remote_url.split("/")[-1]:
        maybe_ext = "." + remote_url.split("/")[-1].split(".")[-1].split("?")[0]
        if len(maybe_ext) <= 5:
            ext = maybe_ext

    filename = f"{kind}_{uuid.uuid4().hex}{ext}"
    local_path = os.path.join(settings.ASSET_STORAGE_DIR, filename)

    resp = requests.get(remote_url, timeout=120, stream=True, headers=headers or {})
    resp.raise_for_status()
    with open(local_path, "wb") as f:
        for chunk in resp.iter_content(chunk_size=8192):
            f.write(chunk)

    return {"url": f"/{local_path}".replace("\\", "/"), "filename": filename}


# --------------------------------------------------------------------------- #
# Parse the <presentation>/<poster> markup the agent emits (same regex logic
# mirrored from the frontend so exports match what's rendered on screen)
# --------------------------------------------------------------------------- #

def parse_presentation_markup(content: str):
    match = re.search(r"<presentation([\s\S]*?)>([\s\S]*?)</presentation>", content, re.IGNORECASE)
    if not match:
        return None
    attrs, body = match.group(1), match.group(2)
    title_match = re.search(r'title=["\']([^"\']+)["\']', attrs, re.IGNORECASE)
    title = title_match.group(1) if title_match else "Presentation"

    slides = []
    for smatch in re.finditer(r"<slide([\s\S]*?)>([\s\S]*?)</slide>", body, re.IGNORECASE):
        sattrs, scontent = smatch.group(1), smatch.group(2)
        stitle_match = re.search(r'title=["\']([^"\']+)["\']', sattrs, re.IGNORECASE)
        stitle = stitle_match.group(1) if stitle_match else "Slide"
        slides.append({"title": stitle, "content": scontent.strip()})

    return {"title": title, "slides": slides}


def parse_poster_markup(content: str):
    match = re.search(r"<poster([\s\S]*?)>([\s\S]*?)</poster>", content, re.IGNORECASE)
    if not match:
        return None
    attrs, body = match.group(1), match.group(2)

    def get_attr(name, default):
        m = re.search(rf'{name}=["\']([^"\']+)["\']', attrs, re.IGNORECASE)
        return m.group(1) if m else default

    title = get_attr("title", "Research Poster")
    authors = get_attr("authors", "")
    domain = get_attr("domain", "")

    sections = []
    for smatch in re.finditer(r"<section([\s\S]*?)>([\s\S]*?)</section>", body, re.IGNORECASE):
        sattrs, scontent = smatch.group(1), smatch.group(2)
        stitle_match = re.search(r'title=["\']([^"\']+)["\']', sattrs, re.IGNORECASE)
        stitle = stitle_match.group(1) if stitle_match else "Section"
        sections.append({"title": stitle, "content": scontent.strip()})

    return {"title": title, "authors": authors, "domain": domain, "sections": sections}


def _strip_markdown_images(text: str):
    """Extract (clean_text, [image_urls]) - pulls out ![alt](url) so slide
    text and images can be placed separately in the pptx/pdf."""
    urls = re.findall(r"!\[[^\]]*\]\(([^)]+)\)", text)
    clean = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    clean = re.sub(r"^\s*-\s+", "• ", clean, flags=re.MULTILINE)
    clean = re.sub(r"[#*`]", "", clean).strip()
    return clean, urls


def _download_to_tmp(url: str) -> str:
    local = os.path.join(settings.EXPORT_STORAGE_DIR, f"tmp_{uuid.uuid4().hex}.png")
    try:
        r = requests.get(url, timeout=60)
        r.raise_for_status()
        with open(local, "wb") as f:
            f.write(r.content)
        return local
    except Exception as e:
        logger.warning(f"Could not download image for export ({url}): {e}")
        return ""


# --------------------------------------------------------------------------- #
# Real .pptx builder
# --------------------------------------------------------------------------- #

def build_pptx(title: str, slides: list) -> str:
    """Builds a real .pptx file from parsed slide data. Returns local file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Generated by Atlas"

    for slide_data in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data["title"]

        clean_text, image_urls = _strip_markdown_images(slide_data["content"])
        body = slide.placeholders[1].text_frame
        body.word_wrap = True
        lines = [l.strip() for l in clean_text.split("\n") if l.strip()]
        if lines:
            body.text = lines[0]
            for line in lines[1:]:
                p = body.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
        else:
            body.text = ""

        for img_url in image_urls[:1]:
            local_img = _download_to_tmp(img_url)
            if local_img:
                try:
                    slide.shapes.add_picture(local_img, Inches(5.2), Inches(1.5), width=Inches(4))
                except Exception as e:
                    logger.warning(f"Failed to embed image in pptx: {e}")

    out_path = os.path.join(settings.EXPORT_STORAGE_DIR, f"presentation_{uuid.uuid4().hex}.pptx")
    prs.save(out_path)
    return out_path


# --------------------------------------------------------------------------- #
# Real poster PDF builder (reportlab if present, else falls back to the
# pure-Python generator already used for text exports)
# --------------------------------------------------------------------------- #

def build_poster_pdf(title: str, authors: str, domain: str, sections: list) -> str:
    out_path = os.path.join(settings.EXPORT_STORAGE_DIR, f"poster_{uuid.uuid4().hex}.pdf")
    try:
        from reportlab.lib.pagesizes import landscape, A2
        from reportlab.lib.units import cm
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader

        page_size = landscape(A2)
        c = canvas.Canvas(out_path, pagesize=page_size)
        width, height = page_size

        c.setFillColorRGB(0.19, 0.15, 0.1)
        c.setFont("Helvetica-Bold", 30)
        c.drawCentredString(width / 2, height - 3 * cm, title[:90])
        c.setFont("Helvetica-Oblique", 16)
        subtitle = f"{authors}  |  {domain}" if authors or domain else ""
        if subtitle:
            c.drawCentredString(width / 2, height - 4 * cm, subtitle)

        col_width = (width - 4 * cm) / 3
        col_x = [1.5 * cm, 1.5 * cm + col_width + 0.5 * cm, 1.5 * cm + 2 * (col_width + 0.5 * cm)]
        col_y = [height - 6 * cm] * 3

        for i, sec in enumerate(sections):
            ci = i % 3
            clean_text, image_urls = _strip_markdown_images(sec["content"])

            c.setFillColorRGB(0.54, 0.44, 0.19)
            c.setFont("Helvetica-Bold", 14)
            c.drawString(col_x[ci], col_y[ci], sec["title"][:60])
            col_y[ci] -= 0.8 * cm

            c.setFillColorRGB(0.15, 0.12, 0.1)
            c.setFont("Helvetica", 10)
            for line in _wrap_text(clean_text, 55):
                if col_y[ci] < 2 * cm:
                    break
                c.drawString(col_x[ci], col_y[ci], line)
                col_y[ci] -= 0.45 * cm

            for img_url in image_urls[:1]:
                local_img = _download_to_tmp(img_url)
                if local_img:
                    try:
                        img = ImageReader(local_img)
                        c.drawImage(img, col_x[ci], col_y[ci] - 4 * cm, width=col_width, height=4 * cm,
                                    preserveAspectRatio=True, anchor='n')
                        col_y[ci] -= 4.3 * cm
                    except Exception as e:
                        logger.warning(f"Failed to embed image in poster pdf: {e}")

            col_y[ci] -= 0.5 * cm

        c.save()
        return out_path
    except ImportError:
        logger.warning("reportlab not installed; falling back to plain-text PDF for poster export.")
        from main import generate_simple_pdf  # pure-python fallback already in main.py
        lines = [f"{title}", f"{authors} | {domain}", ""]
        for sec in sections:
            lines.append(f"== {sec['title']} ==")
            clean_text, _ = _strip_markdown_images(sec["content"])
            lines.append(clean_text)
            lines.append("")
        pdf_bytes = generate_simple_pdf(title, lines)
        with open(out_path, "wb") as f:
            f.write(pdf_bytes)
        return out_path


def _wrap_text(text: str, width: int):
    import textwrap
    wrapped = []
    for paragraph in text.split("\n"):
        if not paragraph.strip():
            wrapped.append("")
            continue
        wrapped.extend(textwrap.wrap(paragraph, width=width) or [""])
    return wrapped